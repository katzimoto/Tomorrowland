"""PPTX text extractor."""

from __future__ import annotations

import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from services.extraction.base import ExtractionResult, LocationSegment


class PptxExtractor:
    """Extract text from PowerPoint .pptx files using python-pptx."""

    def extract(self, path: Path) -> ExtractionResult:
        """Return concatenated text from all slides with slide-number segments."""
        try:
            prs = Presentation(str(path))
            slide_texts: list[str] = []
            segments: list[LocationSegment] = []
            offset = 0
            for i, slide in enumerate(prs.slides, 1):
                slide_parts: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_parts.append(shape.text)
                if slide_parts:
                    slide_text = "\n".join(slide_parts)
                    slide_texts.append(slide_text)
                    end = offset + len(slide_text)
                    segments.append(
                        LocationSegment(
                            start_char=offset,
                            end_char=end,
                            page_number=i,
                        )
                    )
                    offset = end + 1  # +1 for the newline separator between slides
            return ExtractionResult(
                text="\n".join(slide_texts),
                location_segments=segments,
            )
        except (OSError, KeyError, ValueError, zipfile.BadZipFile, PackageNotFoundError):
            return ExtractionResult(text="")
