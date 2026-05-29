"""Markdown-structured Office extractors for improved RAG chunking.

Enabled for OOXML MIME types (DOCX, PPTX, XLSX) when ``ENABLE_MARKITDOWN=true``
is set.  Implemented natively with python-docx, python-pptx, and openpyxl —
the same libraries the markitdown package uses internally — because markitdown
0.1.x requires ``magika<0.7`` which conflicts with this project's ``magika>=1.0``
dependency.

Each converter falls back to the original extractor when Markdown output is
empty or an error occurs, so extraction always degrades gracefully.

Output is for RAG/search chunking only — not for preview rendering.
"""

from __future__ import annotations

import logging
from collections.abc import Callable
from pathlib import Path

from services.extraction.base import ExtractionResult, Extractor

logger = logging.getLogger(__name__)


# ---------------------------------------------------------------------------
# Per-format Markdown converters
# ---------------------------------------------------------------------------


def _docx_to_markdown(path: Path) -> str:
    """Convert a DOCX file to Markdown preserving headings, lists, and tables."""
    import zipfile

    from docx import Document
    from docx.opc.exceptions import PackageNotFoundError
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    try:
        doc = Document(str(path))
    except (OSError, KeyError, ValueError, zipfile.BadZipFile, PackageNotFoundError):
        return ""

    parts: list[str] = []

    def _table_to_md(table: Table) -> str:
        rows: list[list[str]] = []
        for row in table.rows:
            seen_tc: set[int] = set()
            cells: list[str] = []
            for cell in row.cells:
                tc_id = id(cell._tc)  # noqa: SLF001
                if tc_id in seen_tc:
                    continue
                seen_tc.add(tc_id)
                cells.append(cell.text.strip().replace("\n", " "))
            rows.append(cells)
        if not rows:
            return ""
        col_count = max(len(r) for r in rows)
        # Pad all rows to the same width.
        rows = [r + [""] * (col_count - len(r)) for r in rows]
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
        return "\n".join(filter(None, [header, separator, body]))

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag
        if tag == "p":
            para = Paragraph(block, doc)
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name if para.style else ""
            style_lower = style.lower()
            if style_lower.startswith("heading"):
                # "Heading 1" → level 1, "Heading 2" → level 2, etc.
                parts_style = style.split()
                level = int(parts_style[-1]) if parts_style and parts_style[-1].isdigit() else 1
                level = min(level, 6)
                parts.append("#" * level + " " + text)
            elif style_lower in ("list bullet", "list bullet 2", "list paragraph"):
                parts.append("- " + text)
            elif style_lower.startswith("list number"):
                parts.append("1. " + text)
            else:
                parts.append(text)
        elif tag == "tbl":
            tbl = Table(block, doc)
            md_table = _table_to_md(tbl)
            if md_table:
                parts.append(md_table)

    return "\n\n".join(parts)


def _pptx_to_markdown(path: Path) -> str:
    """Convert a PPTX file to Markdown preserving slide titles and bullets."""
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401 — imported to confirm pptx is available

    try:
        prs = Presentation(str(path))
    except Exception:
        return ""

    parts: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        title_text = ""
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title_text = slide.shapes.title.text_frame.text.strip()
        if title_text:
            slide_parts.append(f"## {title_text}")
        else:
            slide_parts.append(f"## Slide {slide_num}")

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                level = para.level or 0
                bullet = "  " * level + "- "
                slide_parts.append(bullet + text)

        if len(slide_parts) > 1:
            parts.append("\n".join(slide_parts))

    return "\n\n".join(parts)


def _xlsx_to_markdown(path: Path) -> str:
    """Convert an XLSX file to Markdown with sheet headings and table rows."""
    import openpyxl

    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    except Exception:
        return ""

    parts: list[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(cells)
        if not rows:
            continue

        col_count = max(len(r) for r in rows)
        rows = [r + [""] * (col_count - len(r)) for r in rows]
        header = "| " + " | ".join(rows[0]) + " |"
        separator = "| " + " | ".join(["---"] * col_count) + " |"
        body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
        sheet_md = f"## {sheet_name}\n\n" + "\n".join(filter(None, [header, separator, body]))
        parts.append(sheet_md)

    wb.close()
    return "\n\n".join(parts)


# ---------------------------------------------------------------------------
# Extractor class
# ---------------------------------------------------------------------------


class MarkItDownExtractor:
    """Wrap a format-specific Markdown converter with a fallback extractor.

    When the Markdown output is empty (e.g. encrypted or unsupported variant),
    *fallback* is called so the original extractor remains the safety net.
    """

    def __init__(
        self,
        convert: Callable[[Path], str],
        fallback: Extractor,
    ) -> None:
        self._convert = convert
        self._fallback = fallback

    def extract(self, path: Path) -> ExtractionResult:
        if not path.exists():
            return ExtractionResult(text="")
        try:
            text = self._convert(path).strip()
            if text:
                return ExtractionResult(text=text)
        except Exception:
            logger.debug("markdown conversion failed path=%s", path, exc_info=True)
        return self._fallback.extract(path)
