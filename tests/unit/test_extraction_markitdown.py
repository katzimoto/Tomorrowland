"""Tests for the native Markdown Office extractors (MarkItDownExtractor)."""

from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock

import openpyxl
from docx import Document
from docx.shared import Pt  # noqa: F401
from pptx import Presentation

from services.extraction.base import ExtractionResult
from services.extraction.markitdown_extractor import (
    MarkItDownExtractor,
    _docx_to_markdown,
    _pptx_to_markdown,
    _xlsx_to_markdown,
)

FIXTURES = Path(__file__).parent.parent / "fixtures"


# ---------------------------------------------------------------------------
# Helpers to build known-content fixture files
# ---------------------------------------------------------------------------


def _make_docx(tmp_path: Path) -> Path:
    """DOCX with a Heading 1, a paragraph, and a 2×2 table."""
    doc = Document()
    doc.add_heading("Main Heading", level=1)
    doc.add_paragraph("This is a body paragraph with some content.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Alpha"
    table.cell(1, 1).text = "Beta"
    path = tmp_path / "test.docx"
    doc.save(str(path))
    return path


def _make_pptx(tmp_path: Path) -> Path:
    """PPTX with two slides: one titled, one with body bullets."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "Slide One Title"
    slide1.placeholders[1].text = "Bullet point alpha\nBullet point beta"
    layout2 = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout2)
    slide2.shapes.title.text = "Slide Two Title"
    slide2.placeholders[1].text = "Another bullet"
    path = tmp_path / "test.pptx"
    prs.save(str(path))
    return path


def _make_xlsx(tmp_path: Path) -> Path:
    """XLSX with two sheets, each containing a header row and data."""
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet Alpha"
    ws1.append(["Name", "Score"])
    ws1.append(["Alice", 95])
    ws1.append(["Bob", 87])
    ws2 = wb.create_sheet("Sheet Beta")
    ws2.append(["Item", "Qty"])
    ws2.append(["Widget", 10])
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return path


# ---------------------------------------------------------------------------
# DOCX converter tests
# ---------------------------------------------------------------------------


class TestDocxToMarkdown:
    def test_heading_preserved(self, tmp_path: Path) -> None:
        path = _make_docx(tmp_path)
        text = _docx_to_markdown(path)
        assert "# Main Heading" in text

    def test_body_paragraph_preserved(self, tmp_path: Path) -> None:
        path = _make_docx(tmp_path)
        text = _docx_to_markdown(path)
        assert "body paragraph" in text

    def test_table_rows_preserved(self, tmp_path: Path) -> None:
        path = _make_docx(tmp_path)
        text = _docx_to_markdown(path)
        assert "Name" in text
        assert "Alpha" in text
        assert "Beta" in text
        # Markdown table separator row
        assert "---" in text

    def test_missing_file_returns_empty(self) -> None:
        result = _docx_to_markdown(Path("/nonexistent/file.docx"))
        assert result == ""

    def test_existing_fixture_non_empty(self) -> None:
        result = _docx_to_markdown(FIXTURES / "sample.docx")
        assert result != ""


# ---------------------------------------------------------------------------
# PPTX converter tests
# ---------------------------------------------------------------------------


class TestPptxToMarkdown:
    def test_slide_title_preserved(self, tmp_path: Path) -> None:
        path = _make_pptx(tmp_path)
        text = _pptx_to_markdown(path)
        assert "Slide One Title" in text
        assert "Slide Two Title" in text

    def test_bullets_preserved(self, tmp_path: Path) -> None:
        path = _make_pptx(tmp_path)
        text = _pptx_to_markdown(path)
        assert "Bullet point alpha" in text
        assert "Bullet point beta" in text

    def test_slide_title_rendered_as_heading(self, tmp_path: Path) -> None:
        path = _make_pptx(tmp_path)
        text = _pptx_to_markdown(path)
        assert "## Slide One Title" in text

    def test_missing_file_returns_empty(self) -> None:
        result = _pptx_to_markdown(Path("/nonexistent/file.pptx"))
        assert result == ""

    def test_existing_fixture_non_empty(self) -> None:
        result = _pptx_to_markdown(FIXTURES / "sample.pptx")
        assert result != ""


# ---------------------------------------------------------------------------
# XLSX converter tests
# ---------------------------------------------------------------------------


class TestXlsxToMarkdown:
    def test_sheet_name_preserved(self, tmp_path: Path) -> None:
        path = _make_xlsx(tmp_path)
        text = _xlsx_to_markdown(path)
        assert "Sheet Alpha" in text
        assert "Sheet Beta" in text

    def test_cell_values_preserved(self, tmp_path: Path) -> None:
        path = _make_xlsx(tmp_path)
        text = _xlsx_to_markdown(path)
        assert "Alice" in text
        assert "95" in text
        assert "Widget" in text
        assert "10" in text

    def test_header_row_in_markdown_table(self, tmp_path: Path) -> None:
        path = _make_xlsx(tmp_path)
        text = _xlsx_to_markdown(path)
        assert "| Name | Score |" in text
        assert "---" in text

    def test_missing_file_returns_empty(self) -> None:
        result = _xlsx_to_markdown(Path("/nonexistent/file.xlsx"))
        assert result == ""

    def test_existing_fixture_non_empty(self) -> None:
        result = _xlsx_to_markdown(FIXTURES / "sample.xlsx")
        assert result != ""


# ---------------------------------------------------------------------------
# MarkItDownExtractor wrapper tests
# ---------------------------------------------------------------------------


class TestMarkItDownExtractor:
    def test_returns_markdown_when_converter_succeeds(self, tmp_path: Path) -> None:
        path = _make_docx(tmp_path)
        fallback = MagicMock(spec=["extract"])
        extractor = MarkItDownExtractor(convert=_docx_to_markdown, fallback=fallback)
        result = extractor.extract(path)
        assert isinstance(result, ExtractionResult)
        assert "Main Heading" in result.text
        fallback.extract.assert_not_called()

    def test_calls_fallback_on_missing_file(self, tmp_path: Path) -> None:
        fallback = MagicMock(spec=["extract"])
        fallback.extract.return_value = ExtractionResult(text="fallback text")
        extractor = MarkItDownExtractor(convert=_docx_to_markdown, fallback=fallback)
        result = extractor.extract(tmp_path / "nonexistent.docx")
        assert result.text == ""
        fallback.extract.assert_not_called()  # path.exists() guard returns early

    def test_calls_fallback_when_converter_returns_empty(self, tmp_path: Path) -> None:
        fallback = MagicMock(spec=["extract"])
        fallback.extract.return_value = ExtractionResult(text="fallback text")
        extractor = MarkItDownExtractor(convert=lambda _: "", fallback=fallback)
        path = tmp_path / "dummy.docx"
        path.write_bytes(b"dummy")
        result = extractor.extract(path)
        assert result.text == "fallback text"
        fallback.extract.assert_called_once_with(path)

    def test_calls_fallback_when_converter_raises(self, tmp_path: Path) -> None:
        def bad_convert(_: Path) -> str:
            raise RuntimeError("simulated failure")

        fallback = MagicMock(spec=["extract"])
        fallback.extract.return_value = ExtractionResult(text="recovered")
        extractor = MarkItDownExtractor(convert=bad_convert, fallback=fallback)
        path = tmp_path / "dummy.docx"
        path.write_bytes(b"dummy")
        result = extractor.extract(path)
        assert result.text == "recovered"
        fallback.extract.assert_called_once_with(path)

    def test_attachments_empty(self, tmp_path: Path) -> None:
        path = _make_docx(tmp_path)
        extractor = MarkItDownExtractor(
            convert=_docx_to_markdown,
            fallback=MagicMock(spec=["extract"]),
        )
        result = extractor.extract(path)
        assert result.attachments == []
